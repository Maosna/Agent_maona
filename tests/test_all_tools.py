#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Maona 全工具 + 全技能模拟测试 - F:\工具\测试\maona_full_test"""

import asyncio
import json
import os
import sys
import time
from datetime import datetime
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent / "backend"))
os.chdir(str(Path(__file__).resolve().parent))

TEST_DIR = Path("F:/工具/测试/maona_full_test")
TEST_DIR.mkdir(parents=True, exist_ok=True)
RESULT_FILE = TEST_DIR / f"result_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"

results = []
pc, fc, sc = 0, 0, 0


def log(name, ok, detail="", skip=False):
    global pc, fc, sc
    tag = "OK" if ok else ("SKIP" if skip else "FAIL")
    if skip: sc += 1
    elif ok: pc += 1
    else: fc += 1
    s = str(detail)[:100]
    print(f"  [{tag}] {name}: {s}")
    results.append({"name": name, "ok": ok, "detail": s, "skip": skip})


async def call(name, args):
    from tools.dispatcher import execute_tool
    return await execute_tool(name, args)


def check_in(text):
    return lambda r: text in r

def check_len(n):
    return lambda r: len(r) > n

def check_ok():
    return lambda r: "error" not in r.lower()[:30]

# ====== 依赖检查 ======
has_pil = False
has_docx = False
has_xlsx = False
has_pdf = False
has_pptx = False
has_playwright = False

try:
    from PIL import Image; has_pil = True
except: pass

try:
    from docx import Document; has_docx = True
except: pass

try:
    import openpyxl; has_xlsx = True
except: pass

try:
    from reportlab.pdfgen import canvas; has_pdf = True
except: pass

try:
    from pptx import Presentation; has_pptx = True
except: pass

try:
    import subprocess
    r = subprocess.run(["playwright", "--version"], capture_output=True, text=True)
    has_playwright = r.returncode == 0
except: pass

# ====== 准备测试文件 ======
def prep_files():
    # txt
    (TEST_DIR / "test.txt").write_text("Hello Maona!", encoding="utf-8")
    # csv
    (TEST_DIR / "test.csv").write_text("col1,col2\na,1\nb,2\n", encoding="utf-8")
    # docx
    if has_docx:
        doc = Document()
        doc.add_heading("Test", 0)
        doc.add_paragraph("test content")
        doc.save(str(TEST_DIR / "test.docx"))
    # xlsx
    if has_xlsx:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws['A1'] = "Name"; ws['A2'] = "Test"
        wb.save(str(TEST_DIR / "test.xlsx"))
    # pdf
    if has_pdf:
        c = canvas.Canvas(str(TEST_DIR / "test.pdf"))
        c.drawString(100, 750, "Test PDF")
        c.save()
    # pptx
    if has_pptx:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"
        prs.save(str(TEST_DIR / "test.pptx"))
    # image
    if has_pil:
        Image.new('RGB', (100, 100), color='blue').save(str(TEST_DIR / "test.png"))
    # html
    (TEST_DIR / "index.html").write_text("<h1>Test</h1>", encoding="utf-8")


async def run_tests():
    global pc, fc, sc

    # ====== 1. 文件操作 ======
    print("\n=== 1. 文件操作 ===")
    await test("write_file", "write_file", {"path": str(TEST_DIR/"w.txt"), "content": "Hello"}, check_ok())
    await test("read_file", "read_file", {"path": str(TEST_DIR/"w.txt")}, check_in("Hello"))
    await test("list_files", "list_files", {"path": str(TEST_DIR)}, check_in("w.txt"))
    await test("search_content", "search_content", {"pattern": "Hello", "path": str(TEST_DIR)}, check_in("w.txt"))
    await test("edit_file", "edit_file", {"path": str(TEST_DIR/"w.txt"), "old_string": "Hello", "new_string": "Hi"}, check_ok())
    await test("rename_file", "rename_file", {"path": str(TEST_DIR/"w.txt"), "new_name": str(TEST_DIR/"w2.txt")}, check_ok())
    await test("delete_file", "delete_file", {"path": str(TEST_DIR/"w2.txt")}, check_ok())

    # ====== 2. 命令执行 ======
    print("\n=== 2. 命令执行 ===")
    await test("run_command", "run_command", {"command": "echo ok123"}, check_in("ok123"))
    await test("run_python", "run_python", {"code": "print('py_ok')"}, check_in("py_ok"))

    # ====== 3. 网络 ======
    print("\n=== 3. 网络工具 ===")
    await test("web_search", "web_search", {"query": "Python programming"}, check_len(50))

    # ====== 4. 记忆 ======
    print("\n=== 4. 记忆系统 ===")
    ts = int(time.time())
    await test("save_memory", "save_memory", {"content": f"mem_test_{ts}", "category": "test"}, check_ok())
    await test("read_memory", "read_memory", {"query": "test"}, check_len(5))
    await test("save_daily_log", "save_daily_log", {"content": "## Test\n- all tools test"}, check_ok())
    await test("save_bug_fix", "save_bug_fix", {"error_pattern": "test error", "fix_description": "test fix", "file_path": "test.py"}, check_in("已记录"))

    # ====== 5. 文档 ======
    print("\n=== 5. 文档处理 ===")
    await test("read_csv", "read_csv", {"path": str(TEST_DIR/"test.csv"), "n": 3}, check_ok())
    await test("read_docx", "read_docx", {"path": str(TEST_DIR/"test.docx")}, check_ok() if has_docx else None, skip=not has_docx)
    await test("read_xlsx", "read_xlsx", {"path": str(TEST_DIR/"test.xlsx")}, check_ok() if has_xlsx else None, skip=not has_xlsx)
    await test("read_pptx", "read_pptx", {"path": str(TEST_DIR/"test.pptx")}, check_ok() if has_pptx else None, skip=not has_pptx)
    await test("pdf_read", "pdf_read", {"path": str(TEST_DIR/"test.pdf")}, check_ok() if has_pdf else None, skip=not has_pdf)

    # ====== 6. 系统 ======
    print("\n=== 6. 系统工具 ===")
    await test("system_info", "system_info", {}, check_ok())
    await test("count_tokens", "count_tokens", {"text": "Hello World"}, check_ok())
    await test("encode_decode_b64", "encode_decode", {"action": "base64_encode", "text": "hello"}, check_in("aGVsbG8="))
    await test("encode_decode_md5", "encode_decode", {"action": "md5", "text": "hello"}, check_in("5d41402abc4b2a76b9719d911017c592"))
    await test("encode_decode_url", "encode_decode", {"action": "url_encode", "text": "hello world"}, check_ok())
    await test("cost_summary", "cost_summary", {}, check_ok())

    # ====== 7. 图片/媒体 ======
    print("\n=== 7. 图片/媒体 ===")
    await test("image_info", "image_info", {"path": str(TEST_DIR/"test.png")}, check_ok(), skip=not has_pil)
    await test("compress_image", "compress_image", {"path": str(TEST_DIR/"test.png"), "width": 50}, check_ok(), skip=not has_pil)
    await test("text_to_speech", "text_to_speech", {"text": "测试"}, check_ok())

    # ====== 8. 文件工具 ======
    print("\n=== 8. 文件工具 ===")
    await test("zip_create", "zip_archive", {"action": "create", "path": str(TEST_DIR/"test.txt"), "dest": str(TEST_DIR/"t.zip")}, check_ok())
    await test("zip_extract", "zip_archive", {"action": "extract", "path": str(TEST_DIR/"t.zip"), "dest": str(TEST_DIR/"extracted")}, check_ok())
    await test("download_file", "download_file", {"url": "https://httpbin.org/get", "path": str(TEST_DIR/"dl_test.txt")}, check_ok())

    # ====== 9. 任务管理 ======
    print("\n=== 9. 任务管理 ===")
    await test("task_create", "task_create", {"subject": "test_task"}, check_ok())
    await test("task_list", "task_list", {}, check_len(5))

    # ====== 10. 技能 ======
    print("\n=== 10. 技能系统 ===")
    await test("find_skills", "find_skills", {"query": "godot"}, check_len(10))
    await test("load_skill", "load_skill", {"skill_id": "godot-dev"}, check_len(20))

    # ====== 11. 知识库 ======
    print("\n=== 11. 知识库 ===")
    kb = f"kb_test_{int(time.time())}"
    await test("kb_create", "kb_create", {"name": kb}, check_ok())
    await test("kb_add", "kb_add", {"kb": kb, "title": "test", "content": "test content 123"}, check_ok())
    await test("kb_search", "kb_search", {"kb": kb, "query": "test"}, check_ok())

    # ====== 12. 代码索引 ======
    print("\n=== 12. 代码索引 ===")
    bd = str(Path(__file__).parent / "backend")
    await test("project_index", "project_index", {"path": bd}, check_len(20))
    await test("code_search", "code_search", {"query": "dispatcher", "path": bd}, check_ok())

    # ====== 13. LSP ======
    print("\n=== 13. LSP 代码智能 ===")
    (TEST_DIR / "test_lsp.py").write_text("def hello():\n    print('ok')\n", encoding="utf-8")
    await test("lsp_diagnose", "lsp_diagnose", {"filepath": str(TEST_DIR/"test_lsp.py")}, check_ok())
    await test("lsp_outline", "lsp_outline", {"filepath": str(TEST_DIR/"test_lsp.py")}, check_ok())
    await test("lsp_hover", "lsp_hover", {"filepath": str(TEST_DIR/"test_lsp.py"), "line": 1}, check_ok())

    # ====== 14. 部署 ======
    print("\n=== 14. 部署 ===")
    await test("deploy_package", "deploy_package", {"directory": str(TEST_DIR), "output": str(TEST_DIR/"pkg.zip")}, check_ok())
    await test("preview_html", "preview_html", {"content": "<h1>Test</h1>"}, check_ok())

    # ====== 15. 备份 ======
    print("\n=== 15. 备份恢复 ===")
    # 先写一个文件再编辑（应该生成备份）
    bp = str(TEST_DIR / "backup_test.txt")
    await test("backup_write", "write_file", {"path": bp, "content": "v1"}, check_ok())
    await test("backup_edit", "edit_file", {"path": bp, "old_string": "v1", "new_string": "v2"}, check_ok())
    await test("restore_backup", "restore_backup", {"path": bp, "list_only": True}, check_ok())

    # ====== 16. 定时任务 ======
    print("\n=== 16. 定时任务 ===")
    await test("schedule_task", "schedule_task", {"command": "echo t", "interval_minutes": 1440, "description": "test"}, check_ok())
    await test("list_scheduled", "list_scheduled_tasks", {}, check_ok())

    # ====== 17. 技能 CRUD ======
    print("\n=== 17. 技能 CRUD ===")
    sn = f"test_skill_{int(time.time())}"
    await test("skill_create", "skill_create", {"name": sn, "body": "# test skill", "display": "Test Skill", "description": "test"}, check_ok())
    await test("skill_update", "skill_update", {"name": sn, "body": "# test skill v2"}, check_ok())
    await test("skill_delete", "skill_delete", {"name": sn}, check_ok())

    # ====== 18. 环境缓存 ======
    print("\n=== 18. 环境缓存 ===")
    await test("cache_env", "cache_env", {}, check_ok())

    # ====== 19. 对话搜索 ======
    print("\n=== 19. 对话搜索 ===")
    await test("search_conversations", "search_conversations", {"query": "test"}, check_ok())

    # ====== 20. Godot 工具 ======
    print("\n=== 20. Godot 工具 ===")
    await test("check_godot_project", "check_godot_project", {"project_dir": str(TEST_DIR)}, check_ok())
    await test("validate_gdscript", "validate_gdscript", {"project_dir": str(TEST_DIR)}, check_ok())

    # ====== 21. ComfyUI ======
    print("\n=== 21. ComfyUI CLI ===")
    await test("comfy_cli", "comfy_cli", {"command": "version"}, check_ok())

    # ====== 22. 开关/其他 ======
    print("\n=== 22. 开关/其他 ===")
    await test("switch_mode", "switch_mode", {"mode": "craft", "reason": "test"}, check_ok())
    await test("install_pip", "install_pip", {"package": "six"}, check_ok())
    # sql_query (skip - might not be there)
    await test("sql_query", "sql_query", {"path": str(TEST_DIR/"nope.db"), "query": "SELECT 1"}, check_ok(), skip=True)

    # ====== 汇总 ======
    print(f"\n{'='*50}")
    print(f"  PASS: {pc}  FAIL: {fc}  SKIP: {sc}  TOTAL: {pc+fc+sc}")
    print(f"{'='*50}")

    with open(RESULT_FILE, "w", encoding="utf-8") as f:
        json.dump({"pass": pc, "fail": fc, "skip": sc, "total": pc+fc+sc, "results": results}, f, ensure_ascii=False, indent=2)
    print(f"Result: {RESULT_FILE}")


async def test(name, tool, args, check_fn, skip=False):
    if skip:
        log(name, False, f"SKIPPED (缺依赖)", skip=True)
        return
    try:
        r = await call(tool, args)
        ok = check_fn(r)
        log(name, ok, r)
    except Exception as e:
        log(name, False, str(e))


if __name__ == "__main__":
    prep_files()
    asyncio.run(run_tests())

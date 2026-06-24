"""文件操作工具实现"""
import os
import re
from pathlib import Path

# 禁止访问的系统目录
FORBIDDEN_DIRS = [
    "/etc", "/boot", "/sys", "/proc", "/dev",
    "C:\\Windows", "C:\\windows", "C:\\Program Files", "C:\\Program Files (x86)",
    "C:\\Windows\\System32", "C:\\Windows\\SysWOW64",
    "/System", "/Library/System", "/Library",
]

# 敏感文件名（仅在系统目录下禁止）
SENSITIVE_FILES = {".git/config", "id_rsa", "id_ed25519", ".ssh/authorized_keys",
                   "web.config", "appsettings.json", "secrets.yml", "credentials.json"}


def _is_path_safe(p: Path) -> bool:
    """检查路径是否在安全范围内"""
    try:
        resolved = p.expanduser().resolve()
    except Exception:
        return False
    path_str = str(resolved).replace("\\", "/").lower()
    # 检查系统目录
    for forbidden in FORBIDDEN_DIRS:
        fb = forbidden.lower().replace("\\", "/")
        if path_str == fb or path_str.startswith(fb + "/"):
            return False
    # 检查敏感文件
    for sensitive in SENSITIVE_FILES:
        if path_str.endswith(sensitive.lower().replace("\\", "/")):
            return False
    return True


async def read_file(path: str) -> str:
    """读取文件内容"""
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return f"❌ 安全限制：不允许访问系统目录或敏感文件 - {path}"
    if not p.exists():
        return f"错误：文件不存在 - {path}"
    if not p.is_file():
        return f"错误：不是文件 - {path}"
    # 先检查大小，避免大文件撑爆内存
    import os
    try:
        fsize = os.path.getsize(str(p))
        if fsize > 10 * 1024 * 1024:  # 10MB 限制
            return f"错误：文件过大（{fsize / 1024 / 1024:.1f}MB），拒绝读取以防止内存溢出。请用 grep 或 run_command head/tail 分段读取"
    except OSError:
        pass
    try:
        content = p.read_text(encoding="utf-8-sig", errors="replace")
        # Claude Code 风格：不硬截断，让模型上下文窗口做上限
        # 仅超大文件（>100K 字符）标注提示
        if len(content) > 100000:
            content = content[:100000] + f"\n... (文件过大，只读取前 100K 字符。全文共 {len(content)} 字符，请用 grep 或分段读取)"
        return content
    except PermissionError:
        return f"错误：无权限读取 - {path}"
    except Exception as e:
        return f"错误：{e}"


async def search_content(pattern: str, path: str = ".", file_pattern: str = "*") -> str:
    """在目录下搜索包含指定文本的文件（grep）"""
    import fnmatch
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return f"❌ 安全限制：不允许搜索系统目录 - {path}"
    if not p.exists():
        return f"错误：路径不存在 - {path}"
    search_dir = p if p.is_dir() else p.parent
    results = []
    try:
        for root, dirs, files in os.walk(search_dir):
            dirs[:] = [d for d in dirs if not d.startswith(".") and d not in ("node_modules", "__pycache__")]
            for fname in files:
                if not fnmatch.fnmatch(fname, file_pattern):
                    continue
                fpath = os.path.join(root, fname)
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        for i, line in enumerate(f, 1):
                            if pattern in line:
                                results.append(f"{fpath}:{i}: {line.rstrip()[:500]}")
                                if len(results) >= 100:
                                    break
                except Exception:
                    continue
                if len(results) >= 100:
                    break
            if len(results) >= 100:
                break
    except Exception as e:
        return f"搜索错误: {e}"
    if not results:
        return f"在 {search_dir} 中未找到包含「{pattern}」的文件"
    return f"搜索「{pattern}」结果 ({len(results)} 条):\n" + "\n".join(results)


async def edit_file(path: str, old_string: str = "", new_string: str = "",
                     function: str = "", class_name: str = "", method: str = "",
                     replace_body: str = "", append_after: str = "",
                     regex: str = "") -> str:
    """精确替换文件文本，支持结构化编辑（函数/类级别）

    模式（按优先级）：
    1. function=<name> + replace_body=<new_body> → 替换整个函数体
    2. class_name=<name> + method=<name> + replace_body=<new_body> → 替换类方法
    3. regex=<pattern> + new_string → 正则替换（首次匹配）
    4. old_string + new_string → 精确字符串替换（默认）
    5. append_after=<match> + new_string → 在匹配行后追加内容
    """
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return f"❌ 安全限制：不允许编辑系统目录文件 - {path}"
    if not p.exists():
        return f"错误：文件不存在 - {path}"
    try:
        content = p.read_text(encoding="utf-8-sig")
        orig = content
    except Exception as e:
        return f"错误：无法读取文件 - {e}"

    # 模式 1: 函数级替换 (GDScript: func name(params):)
    if function and replace_body:
        pattern = rf"func\s+{re.escape(function)}\s*\([^)]*\)\s*(?:->\s*\w+\s*)?:\s*\n"
        m = re.search(pattern, content)
        if not m:
            return f"错误：找不到函数 func {function}()"
        start = m.end()
        # 找函数体结束（下一个同缩进层的 func 定义）
        rest = content[start:]
        lines = rest.split("\n")
        body_end = len(lines)
        for i, line in enumerate(lines):
            stripped = line.lstrip()
            if stripped and not line.startswith((" ", "\t", "#")) and not line.strip().startswith("#"):
                if re.match(rf"func\s+\w+\s*\(", stripped):
                    body_end = i
                    break
        old_body = "\n".join(lines[:body_end])
        new_content = content[:start] + replace_body + "\n" + "\n".join(lines[body_end:])
        p.write_text(new_content, encoding="utf-8")
        result = f"已替换函数 {function}() 的完整实现"
        validation = _auto_validate(p, new_content)
        return result + validation if validation else result

    # 模式 2: 类方法替换 (GDScript: class_name extends Node: ... func method(): ...)
    if class_name and method and replace_body:
        # 找到类定义
        class_pat = rf"class_name\s+{re.escape(class_name)}\s*\n"
        class_m = re.search(class_pat, content)
        if not class_m:
            # 尝试 class XXX extends YYY:
            class_pat2 = rf"class\s+{re.escape(class_name)}\s*(?:extends\s+\w+)?\s*:\s*\n"
            class_m = re.search(class_pat2, content)
        if not class_m:
            return f"错误：找不到类 {class_name}"
        # 在类内找方法
        class_content = content[class_m.end():]
        method_pat = rf"func\s+{re.escape(method)}\s*\([^)]*\)\s*(?:->\s*\w+\s*)?:\s*\n"
        method_m = re.search(method_pat, class_content)
        if not method_m:
            return f"错误：在类 {class_name} 中找不到方法 {method}()"
        abs_start = class_m.end() + method_m.end()
        # 找方法体结束
        rest = content[abs_start:]
        lines = rest.split("\n")
        body_end = len(lines)
        for i, line in enumerate(lines):
            if i == 0: continue
            stripped = line.strip()
            if stripped and not line.startswith((" ", "\t")):
                body_end = i
                break
        old_body = "\n".join(lines[:body_end])
        new_content = content[:abs_start] + replace_body + "\n" + "\n".join(lines[body_end:])
        p.write_text(new_content, encoding="utf-8")
        result = f"已替换类 {class_name} 的方法 {method}()"
        validation = _auto_validate(p, new_content)
        return result + validation if validation else result

    # 模式 3: 正则替换
    if regex and new_string:
        new_content, n = re.subn(regex, new_string, content, count=1)
        if n == 0:
            return f"错误：正则 {regex} 未匹配到任何内容"
        p.write_text(new_content, encoding="utf-8")
        result = f"已用正则修改 {path}（1 处替换）"
        validation = _auto_validate(p, new_content)
        return result + validation if validation else result

    # 模式 4: 字符串精确替换（原逻辑）
    if old_string and old_string != new_string:
        if old_string == new_string:
            return "错误：old_string 和 new_string 相同，无需修改"
        count = content.count(old_string)
        if count == 0:
            # 智能提示：检查是否是缩进差异
            hint = _find_close_match(old_string, content)
            return f"错误：文件中找不到指定的 old_string。{hint}"
        if count > 1:
            return f"错误：old_string 出现了 {count} 次，不唯一，请提供更多上下文或用 function=/class_name= 参数精确指定"
        new_content = content.replace(old_string, new_string, 1)
        p.write_text(new_content, encoding="utf-8")
        result = f"已修改 {path}"
        if p.suffix.lower() in {".tscn", ".tres", ".res"}:
            result = ("⚠️ 注意：正在编辑 Godot 专有格式文件（{0}），手改极易引入格式错误。"
                      "编辑后请调用 validate_gdscript 验证。\n").format(p.suffix) + result
        validation = _auto_validate(p, new_content)
        return result + validation if validation else result

    # 模式 5: 追加内容
    if append_after and new_string:
        idx = content.find(append_after)
        if idx == -1:
            return f"错误：找不到定位文本 {append_after[:50]}..."
        insert_at = idx + len(append_after)
        # 确保在行尾追加
        next_nl = content.find("\n", insert_at)
        if next_nl != -1:
            insert_at = next_nl
        new_content = content[:insert_at] + "\n" + new_string + content[insert_at:]
        p.write_text(new_content, encoding="utf-8")
        result = f"已在 {append_after[:30]}... 后追加内容到 {path}"
        validation = _auto_validate(p, new_content)
        return result + validation if validation else result

    return "错误：请提供 old_string+new_string、function+replace_body、class_name+method+replace_body 或 regex+new_string"


def _find_close_match(old: str, content: str) -> str:
    """尝试找缩进变体的匹配"""
    lines = old.strip().split("\n")
    if not lines: return "提示：请用 read_file 重新读取文件内容"
    first = lines[0].lstrip()
    for variant in [first, "    " + first, "\t" + first, "  " + first]:
        if variant in content:
            return f"提示：找到相似文本但缩进不同。请复制文件中的精确字符串（含空格/Tab）"
    return "提示：文件可能已被修改，请用 read_file 重新读取当前内容后再编辑"


# === 硬管道验证：写文件后自动校验 ===

def _auto_validate(p: Path, content: str) -> str:
    """写文件后立即校验格式/语法，返回错误提示（无错误返回空串）"""
    suffix = p.suffix.lower()
    hints = []

    # 1. Python 语法检查
    if suffix == ".py":
        try:
            compile(content, str(p), "exec")
        except SyntaxError as e:
            hints.append(f"\n\n## 🔴 写入验证失败\n{p.name} 有语法错误:\n  行 {e.lineno}: {e.msg}\n  {e.text}")
        except Exception:
            pass

    # 2. project.godot 格式检查
    if p.name == "project.godot":
        h = _check_project_godot_content(content)
        if h:
            hints.append(h)

    # 3. GDScript 基础检查
    if suffix == ".gd":
        h = _check_gdscript_basic(content, p.name)
        if h:
            hints.append(h)

    # 4. JSON 语法检查
    if suffix == ".json":
        import json as _json
        try:
            _json.loads(content)
        except _json.JSONDecodeError as e:
            hints.append(f"\n\n## 🔴 JSON 解析失败\n{p.name}: {e}")

    return "".join(hints)


def _check_project_godot_content(content: str) -> str:
    """检查 project.godot 内容"""
    issues = []
    if 'config_version=5' not in content and 'config_version = 5' not in content:
        return ""
    if re.search(r'PACKED_STRING_ARRAY', content):
        issues.append("🔴 PACKED_STRING_ARRAY 大小写错误，应为 PackedStringArray")
    for wrong in ['PACKED_INT32_ARRAY', 'PACKED_FLOAT64_ARRAY']:
        if wrong in content:
            issues.append(f"🔴 {wrong} 大小写错误")
    # 检查 autoload 路径中的 "*res://" 错误（模型常见幻觉）
    star_res = re.findall(r'"\*res://[^"]*"', content)
    if star_res:
        issues.append(f"🔴 autoload 路径有 '*res://' 错误: {', '.join(star_res)}。应改为 res:// 去掉 *")
    if issues:
        return f"\n\n## 🔴 project.godot 格式问题\n" + "\n".join(f"  {i}" for i in issues)
    return ""


def _check_gdscript_basic(content: str, fname: str) -> str:
    """GDScript 基础语法检查"""
    issues = []
    # 检查未闭合的括号
    opens = content.count("(") - content.count(")")
    if opens != 0:
        issues.append(f"🟡 括号不匹配: ( 多 {opens} 个" if opens > 0 else f"🟡 括号不匹配: ) 多 {-opens} 个")
    # 检查 extends 声明
    if not re.search(r'extends\s+\w+', content):
        issues.append("🟡 缺少 extends 声明")
    # 检查缩进一致性（tab vs space）
    if "\t" in content and "    " in content:
        issues.append("🟡 混合使用 Tab 和空格缩进")
    if issues:
        return f"\n\n## 🟡 {fname} 基础检查\n" + "\n".join(f"  {i}" for i in issues)
    return ""


async def write_file(path: str, content: str) -> str:
    """写入文件"""
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return f"❌ 安全限制：不允许写入系统目录或敏感文件 - {path}"
    
    # Godot 专有格式护栏：手写 .tscn 极易产生格式错误（如混入 GDScript 表达式）
    suffix = p.suffix.lower()
    GODOT_SPECIAL = {".tscn", ".tres", ".res"}
    warning = ""
    if suffix in GODOT_SPECIAL:
        warning = ("⚠️ 注意：正在直接写入 Godot 专有格式文件 {0}。手写 {0} 文件极易引入格式错误"
                   "（如 StyleBoxFlat.new() 在 Godot 4 中无效）。\n"
                   "如果当前有可用的 Godot 编辑器（9080 端口监听），强烈建议改用 build_godot_scene。\n"
                   "如果必须手写，写入后应调用 validate_gdscript 验证项目完整性。\n\n").format(suffix)
    
    try:
        old = p.read_text(encoding="utf-8-sig", errors="replace") if p.exists() else None
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content, encoding="utf-8")
        result = warning
        if old is not None:
            try:
                import difflib
                diff = list(difflib.unified_diff(
                    old.splitlines(keepends=True),
                    content.splitlines(keepends=True),
                    fromfile="a/" + str(p.name), tofile="b/" + str(p.name)
                ))
                result += f"已修改文件: {p}\n变更 ({len(diff)//2} 处):\n" + "".join(diff[:30])
            except Exception:
                result += f"已修改文件: {p} ({len(content)} 字符)"
        else:
            result += f"已创建文件: {p} ({len(content)} 字符)"
        # === 硬管道验证：写完后立即自动校验 ===
        validation = _auto_validate(p, content)
        if validation:
            result += validation
        return result
    except PermissionError:
        return f"错误：无权限写入 - {path}"
    except Exception as e:
        return f"错误：{e}"


async def git_diff(path: str, staged: bool = False) -> str:
    """查看 git 仓库的变更（工作区 vs HEAD）"""
    import asyncio
    p = Path(path).expanduser().resolve()
    repo = p if p.is_dir() else p.parent
    try:
        args = ["git", "-C", str(repo), "diff"]
        if staged:
            args.append("--staged")
        proc = await asyncio.create_subprocess_exec(*args, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
        stdout, _ = await proc.communicate()
        out = stdout.decode("utf-8", errors="replace")
        return f"Git 变更 ({repo}):\n{out[:5000]}" if out.strip() else f"Git: 无变更"
    except FileNotFoundError:
        return "错误：未安装 git 或不在 PATH 中"
    except Exception as e:
        return f"git diff 失败: {e}"


async def git_log(path: str, n: int = 10) -> str:
    """查看 git 提交历史"""
    import asyncio
    p = Path(path).expanduser().resolve()
    repo = p if p.is_dir() else p.parent
    try:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", str(repo), "log", f"-{n}", "--oneline", "--decorate",
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, _ = await proc.communicate()
        return f"Git 历史 ({repo}):\n{stdout.decode('utf-8', errors='replace').strip() or '无提交'}"
    except FileNotFoundError:
        return "错误：未安装 git 或不在 PATH 中"
    except Exception as e:
        return f"git log 失败: {e}"


async def list_files(path: str = None) -> str:
    """列出目录内容"""
    if not path:
        path = str(Path.home())
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return f"❌ 安全限制：不允许列出系统目录 - {path}"
    if not p.exists():
        return f"错误：目录不存在 - {path}"
    if not p.is_dir():
        return f"错误：不是目录 - {path}"
    try:
        lines = []
        for item in sorted(p.iterdir()):
            marker = "/" if item.is_dir() else ""
            try:
                size = item.stat().st_size
            except:
                size = 0
            lines.append(f"  {item.name}{marker}  ({size:,} bytes)")
        if not lines:
            return f"目录为空: {p}"
        return f"目录 {p} 的内容:\n" + "\n".join(lines[:200])
    except PermissionError:
        return f"错误：无权限访问 - {path}"
    except Exception as e:
        return f"错误：{e}"

async def pdf_read(path: str) -> str:
    """读取 PDF 文件文本内容"""
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return f"❌ 安全限制：不允许读取系统目录 - {path}"
    if not p.exists():
        return f"错误：文件不存在 - {path}"
    try:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(p))
            text = "\n".join(page.extract_text() or "" for page in reader.pages[:50])
            if text.strip():
                return f"PDF 内容 ({len(reader.pages)} 页, 显示前 50 页):\n{text[:10000]}"
        except ImportError:
            pass
        try:
            import pdfplumber
            with pdfplumber.open(str(p)) as pdf:
                text = "\n".join(page.extract_text() or "" for page in pdf.pages[:50])
                return f"PDF 内容 ({len(pdf.pages)} 页, 显示前 50 页):\n{text[:10000]}"
        except ImportError:
            pass
        return "错误：未安装 PDF 解析库。pip install PyPDF2 pdfplumber"
    except Exception as e:
        return f"PDF 读取失败: {e}"


async def sql_query(path: str, query: str) -> str:
    """查询 SQLite 数据库（仅 SELECT）"""
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return f"❌ 安全限制 - {path}"
    if not p.exists():
        return f"错误：数据库不存在 - {path}"
    import re
    # 严格限制：仅允许以 SELECT 开头（忽略前导空白和注释）的查询
    q_stripped = query.strip()
    if not re.match(r'^\s*--.*?\n|^\s*/\*.*?\*/|^\s*SELECT\b', q_stripped, re.IGNORECASE | re.DOTALL):
        if not re.match(r'^\s*SELECT\b', q_stripped, re.IGNORECASE):
            return f"❌ 安全限制：仅允许 SELECT 查询"
    # 禁止多语句（分号分隔）
    if ';' in q_stripped.rstrip(';'):
        return f"❌ 安全限制：不允许同时执行多条语句"
    try:
        import aiosqlite
        db = await aiosqlite.connect(str(p))
        try:
            cur = await db.execute(query)
            rows = await cur.fetchall()
            if not rows:
                return "查询完成: 0 行"
            cols = [d[0] for d in cur.description] if cur.description else []
            h = " | ".join(cols) if cols else ""
            lines = [f"查询完成: {len(rows)} 行"]
            if h:
                lines.append(h)
                lines.append("-" * min(len(h), 80))
            for r in rows[:50]:
                lines.append(" | ".join(str(v) for v in r))
            if len(rows) > 50:
                lines.append(f"... (还有 {len(rows) - 50} 行)")
            return "\n".join(lines)
        finally:
            await db.close()
    except ImportError:
        return "错误：未安装 aiosqlite"
    except Exception as e:
        return f"SQL 错误: {e}"


async def api_post(url: str, body: str = None, headers: str = None) -> str:
    """发送 HTTP POST 请求"""
    import httpx, json
    from urllib.parse import urlparse
    parsed = urlparse(url)
    if parsed.hostname in ("127.0.0.1", "localhost", "::1"):
        return "❌ 安全限制：禁止 POST 到本地"
    if parsed.hostname.startswith("169.254."):
        return "❌ 安全限制：禁止 POST 云元数据"
    try:
        h = {}
        if headers:
            try:
                h = json.loads(headers)
            except:
                for line in headers.split("\n"):
                    if ":" in line:
                        k, v = line.split(":", 1)
                        h[k.strip()] = v.strip()
        b = body
        try:
            b = json.loads(body)
        except:
            pass
        proxies = None
        for v in ("HTTPS_PROXY", "https_proxy", "HTTP_PROXY", "http_proxy", "ALL_PROXY", "all_proxy"):
            p = os.environ.get(v)
            if p: proxies = p; break
        async with httpx.AsyncClient(timeout=15, follow_redirects=True, proxy=proxies, trust_env=True) as client:
            resp = await client.post(url, json=b if isinstance(b, dict) else None,
                                     content=b if not isinstance(b, dict) else None, headers=h)
            return f"HTTP {resp.status_code}\n{resp.text[:5000]}"
    except Exception as e:
        return f"POST 失败: {e}"


async def rename_file(path: str, new_name: str) -> str:
    """重命名或移动文件"""
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return "❌ 安全限制"
    if not p.exists():
        return f"错误：文件不存在 - {path}"
    new_p = p.parent / new_name if "/" not in new_name and "\\" not in new_name else Path(new_name).expanduser().resolve()
    if not _is_path_safe(new_p):
        return "❌ 安全限制：目标路径不安全"
    try:
        p.rename(new_p)
        return f"✅ 已重命名: {p} -> {new_p}"
    except Exception as e:
        return f"重命名失败: {e}"


async def delete_file(path: str) -> str:
    """安全删除文件（移到回收站而非永久删除）"""
    p = Path(path).expanduser().resolve()
    if not _is_path_safe(p):
        return "❌ 安全限制"
    if not p.exists():
        return f"错误：文件不存在 - {path}"
    try:
        import send2trash
        send2trash.send2trash(str(p))
        return f"✅ 已移到回收站: {p}"
    except ImportError:
        try:
            p.unlink() if p.is_file() else p.rmdir()
            return f"✅ 已删除: {p} (未安装 send2trash，已永久删除)"
        except Exception as e2:
            return f"删除失败: {e2}"
    except Exception as e:
        return f"回收站操作失败: {e}"


async def read_docx(path: str) -> str:
    """读取 Word (.docx) 文件文本"""
    p = Path(path).expanduser().resolve()
    if not p.exists():
        return f"错误：文件不存在 - {path}"
    try:
        from docx import Document
        doc = Document(str(p))
        paras = [para.text for para in doc.paragraphs if para.text.strip()]
        return f"DOCX: {p.name} ({len(paras)} 段落)\n" + "\n".join(paras[:50])
    except ImportError:
        return "错误：未安装 python-docx。pip install python-docx"
    except Exception as e:
        return f"DOCX 读取失败: {e}"


async def read_pptx(path: str) -> str:
    """读取 PPT (.pptx) 幻灯片文本"""
    p = Path(path).expanduser().resolve()
    if not p.exists():
        return f"错误：文件不存在 - {path}"
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append(f"--- 第{i}页 ---\n" + "\n".join(texts))
        return f"PPTX: {p.name} ({len(prs.slides)} 页)\n" + "\n\n".join(slides[:20])
    except ImportError:
        return "错误：未安装 python-pptx。pip install python-pptx"
    except Exception as e:
        return f"PPTX 读取失败: {e}"


async def read_xlsx(path: str, sheet: str = None, n: int = 20) -> str:
    """读取 Excel (.xlsx/.xls) 表格数据"""
    p = Path(path).expanduser().resolve()
    if not p.exists():
        return f"错误：文件不存在 - {path}"
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(p), data_only=True)
        ws = wb[sheet] if sheet else wb.active
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            return "空表格"
        header = [str(c) if c else "" for c in rows[0]]
        lines = [f"XLSX: {p.name} / {ws.title} ({len(rows)-1} 行, {len(header)} 列)"]
        lines.append(" | ".join(header))
        lines.append("-" * 60)
        for row in rows[1:n+1]:
            lines.append(" | ".join(str(c) if c is not None else "" for c in row))
        if len(rows) > n + 1:
            lines.append(f"... (还有 {len(rows) - n - 1} 行)")
        return "\n".join(lines)
    except ImportError:
        return "错误：未安装 openpyxl。pip install openpyxl"
    except Exception as e:
        return f"XLSX 读取失败: {e}"


async def html_preview(html: str, path: str = None) -> str:
    """将 HTML 代码保存为文件并在浏览器中打开预览"""
    import webbrowser, os, asyncio
    dest = Path(path).expanduser().resolve() if path else Path.home() / "Desktop" / "maona_preview.html"
    if not _is_path_safe(dest):
        return f"❌ 安全限制：预览路径不安全 - {dest}"
    dest.parent.mkdir(parents=True, exist_ok=True)
    full_html = f"<!DOCTYPE html>\n<html lang=\"zh\">\n<head>\n<meta charset=\"UTF-8\">\n</head>\n<body>\n{html}\n</body>\n</html>"
    dest.write_text(full_html, encoding="utf-8")
    await asyncio.to_thread(webbrowser.open, f"file://{dest}")
    return f"HTML 预览已打开: {dest}"

